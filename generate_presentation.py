"""
generate_presentation.py — Generate the project presentation (PPTX).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)


def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return tf


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             'Ask Your Data', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
             'Governed AI for Enterprise Data Questions', font_size=28, color=RGBColor(0xAE,0xD6,0xF1),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             'Next Level Challenge Hackathon  |  Team 18', font_size=20, color=RGBColor(0xD5,0xDB,0xDB),
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.6),
             'Collibra Data Intelligence  \u00d7  Databricks Delta Sharing  \u00d7  Anthropic Claude',
             font_size=16, color=RGBColor(0x85,0x92,0x9E), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'The Problem', font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             'AI can generate SQL from natural language \u2014 but without business context, it guesses wrong.',
             font_size=20, color=MID_GRAY)

# Left box - Without context
add_shape_bg(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(0.5),
             '\u274c  Without Governed Context', font_size=22, bold=True, color=RED)

add_bullet_frame(slide, Inches(1.1), Inches(3.2), Inches(5), Inches(3.5), [
    'User asks: "How many active testers in ongoing programs?"',
    'LLM guesses Z_PRJ_STAT = \'ACTIVE\'  (wrong value)',
    'LLM guesses Z_ENG_ST = \'ACTIVE\'  (no such value)',
    'SQL returns 0 results',
    'Answer: "Zero active testers" \u2014 WRONG',
], font_size=15, color=DARK_GRAY)

# Right box - With context
add_shape_bg(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4.5), WHITE)
add_text_box(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.5),
             '\u2705  With Governed Context', font_size=22, bold=True, color=GREEN)

add_bullet_frame(slide, Inches(7.3), Inches(3.2), Inches(5), Inches(3.5), [
    'Same question, same LLM, same data',
    'Collibra defines "Active Tester" = 3 criteria (OR)',
    'Collibra defines "Ongoing" = Z_PRJ_STAT = \'Ongoing\'',
    'SQL uses correct columns & values',
    'Answer: 303 active testers \u2014 CORRECT',
], font_size=15, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Approach
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Our Approach: Two Engines', font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             'Flexibility for exploration + Reliability for governed metrics',
             font_size=20, color=MID_GRAY)

# Engine 1
add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.8), WHITE)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.5),
             'Engine 1: RAG Pipeline (LLM-based)', font_size=22, bold=True, color=ACCENT_BLUE)

add_bullet_frame(slide, Inches(1.1), Inches(3.1), Inches(5), Inches(3.8), [
    '1. Ingest all governed context from Collibra API',
    '2. Inject definitions into LLM system prompt',
    '3. User asks a question in plain English',
    '4. LLM generates SQL grounded in definitions',
    '5. Execute SQL against Databricks data',
    '6. Return answer with full traceability',
    '',
    'Best for: Open-ended, exploratory questions',
], font_size=14)

# Engine 2
add_shape_bg(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.8), WHITE)
add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(0.5),
             'Engine 2: Deterministic Rulebook', font_size=22, bold=True, color=GREEN)

add_bullet_frame(slide, Inches(7.3), Inches(3.1), Inches(5), Inches(3.8), [
    '1. Extract governed definitions from Collibra',
    '2. Verify columns & values against real data',
    '3. Encode logic as pure Python/pandas code',
    '4. No LLM in the loop \u2014 zero variability',
    '5. Returns answer + full provenance',
    '6. Same answer every time (20/20 verified)',
    '',
    'Best for: Mission-critical governed metrics',
], font_size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture (4 Layers)
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Architecture: Collibra\u2019s Four-Layer Model', font_size=36, bold=True, color=DARK_BLUE)

layers = [
    ("Layer 1: Business Vocabulary", "Glossary + Metrics",
     '"Active Tester" = 3-criteria OR definition\n"Active Tester Flag" = computed measure',
     RGBColor(0x2E, 0x86, 0xC1)),
    ("Layer 2: Logical Layer", "Data Entities & Attributes",
     'Maps business concepts to logical data structures\n"Completed Activities" \u2192 logical attribute',
     RGBColor(0x28, 0xB4, 0x63)),
    ("Layer 3: Physical Layer", "Schemas, Tables, Columns",
     'Z_ACT_CMP = "Completed Activities"\nZ_PRJ_STAT = "Project Status"',
     RGBColor(0xF3, 0x9C, 0x12)),
    ("Layer 4: Data Layer", "Databricks via Delta Sharing",
     '8 tables, ~20K total rows\nzcc_prt_mtrc, zcc_tkt_itm, zcc_prj_hdr, ...',
     RGBColor(0xE7, 0x4C, 0x3C)),
]

y_pos = Inches(1.5)
for title, subtitle, desc, color in layers:
    add_shape_bg(slide, Inches(0.8), y_pos, Inches(0.15), Inches(1.2), color)
    add_text_box(slide, Inches(1.2), y_pos, Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=DARK_BLUE)
    add_text_box(slide, Inches(1.2), y_pos + Inches(0.4), Inches(4), Inches(0.3),
                 subtitle, font_size=14, color=MID_GRAY)
    add_text_box(slide, Inches(5.5), y_pos + Inches(0.1), Inches(7), Inches(1),
                 desc, font_size=14, color=DARK_GRAY)
    y_pos += Inches(1.4)

add_text_box(slide, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.5),
             'Traceability chain:  Business Question \u2192 Business Term \u2192 Logical Attribute \u2192 Physical Column \u2192 Data',
             font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Context Ingestion
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Context Ingestion from Collibra', font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             'collibra_client.py extracts the complete governed context via REST API',
             font_size=18, color=MID_GRAY)

stats = [
    ("52", "Business Terms", "Official definitions from\nthe Glossary domain"),
    ("4", "Measures", "Calculation rules for\ngoverned metrics"),
    ("230", "Relations", "Links connecting assets\nacross all 4 layers"),
    ("65", "Column Mappings", "Physical column \u2192\nBusiness meaning"),
]

x_pos = Inches(0.8)
for num, label, desc in stats:
    add_shape_bg(slide, x_pos, Inches(2.2), Inches(2.7), Inches(2.5), WHITE)
    add_text_box(slide, x_pos, Inches(2.4), Inches(2.7), Inches(1),
                 num, font_size=48, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x_pos, Inches(3.3), Inches(2.7), Inches(0.4),
                 label, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x_pos, Inches(3.8), Inches(2.7), Inches(0.8),
                 desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(3.1)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8),
             'Key insight: Collibra stores definitions as HTML. We strip tags with BeautifulSoup.\n'
             'Bug found & fixed: numeric attribute values (e.g., 255.0) caused TypeError \u2014\n'
             'added type check before HTML parsing.\n\n'
             'Output: governed_context.json \u2014 the single source of truth for all downstream components.',
             font_size=14, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Active Tester Definition
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Governed Definition: Active Tester', font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             'From Collibra Business Term (asset 0198c234-11fe-73ff-be9b-c91312850031)',
             font_size=16, color=MID_GRAY)

# Three criteria boxes
criteria = [
    ("Criterion 1", "Feedback Tickets \u2265 3",
     "COUNT(DISTINCT Z_TKT_NR)\nFROM zcc_tkt_itm\nGROUP BY Z_SMTP_ADR",
     "73 participants met this", ACCENT_BLUE),
    ("Criterion 2", "Surveys Completed \u2265 2",
     "SUM(Z_SRV_CMP)\nFROM zcc_prt_mtrc\nGROUP BY Z_SMTP_ADR",
     "284 participants met this", GREEN),
    ("Criterion 3", "Completed > Incomplete",
     "Z_ACT_CMP > Z_ACT_INC\n+ Z_ACT_BLK + Z_ACT_OPT\nGROUP BY Z_SMTP_ADR",
     "64 participants met this", ORANGE),
]

x_pos = Inches(0.8)
for title, subtitle, sql, result, color in criteria:
    add_shape_bg(slide, x_pos, Inches(2), Inches(3.7), Inches(3.8), WHITE)
    add_shape_bg(slide, x_pos, Inches(2), Inches(3.7), Inches(0.1), color)
    add_text_box(slide, x_pos + Inches(0.2), Inches(2.3), Inches(3.3), Inches(0.4),
                 title, font_size=16, bold=True, color=MID_GRAY)
    add_text_box(slide, x_pos + Inches(0.2), Inches(2.7), Inches(3.3), Inches(0.5),
                 subtitle, font_size=20, bold=True, color=DARK_BLUE)
    add_text_box(slide, x_pos + Inches(0.2), Inches(3.4), Inches(3.3), Inches(1.2),
                 sql, font_size=12, color=DARK_GRAY, font_name="Consolas")
    add_text_box(slide, x_pos + Inches(0.2), Inches(4.8), Inches(3.3), Inches(0.5),
                 result, font_size=14, bold=True, color=color)
    x_pos += Inches(4.1)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             'Combination: OR logic (any ONE criterion met = active tester)\n'
             'Scope: Only "Ongoing" projects (Z_PRJ_STAT = \'Ongoing\', 20 of 100 projects)\n'
             'Identifier: Z_SMTP_ADR (hashed email) \u2014 unique across tables',
             font_size=15, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Rulebook
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'The Rulebook: From Definition to Code', font_size=36, bold=True, color=DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             'Why: LLMs can misinterpret definitions. For governed metrics, encode logic directly.',
             font_size=18, color=MID_GRAY)

steps = [
    ("1", "Extract", "Read the governed\ndefinition from\nCollibra API"),
    ("2", "Identify", "Map criteria to\nphysical columns\nvia relations"),
    ("3", "Verify", "Confirm values\nagainst actual data\n(\'Ongoing\' not \'Active\')"),
    ("4", "Encode", "Write pure pandas\ncode implementing\nthe exact logic"),
    ("5", "Provenance", "Return answer +\nCollibra asset IDs\n+ full audit trail"),
]

x_pos = Inches(0.5)
for num, title, desc in steps:
    add_shape_bg(slide, x_pos, Inches(2.2), Inches(2.2), Inches(2.8), WHITE)
    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.75), Inches(2.4), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x_pos + Inches(0.1), Inches(3.3), Inches(2), Inches(0.4),
                 title, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x_pos + Inches(0.1), Inches(3.7), Inches(2), Inches(1),
                 desc, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(2.5)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             'Result: deterministic_engine.py produces identical answers on every run.\n'
             'No LLM variability. No prompt sensitivity. No hallucinated filter values.\n'
             'Every line of code traces back to a Collibra asset ID.',
             font_size=16, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Results
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Results', font_size=36, bold=True, color=DARK_BLUE)

# Big number
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(3.5), WHITE)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(4.5), Inches(1.2),
             '303', font_size=96, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(3.2), Inches(4.5), Inches(0.5),
             'Active Testers in Ongoing Programs', font_size=18, bold=True, color=DARK_BLUE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4.5), Inches(0.5),
             '20/20 runs identical \u2014 zero variance', font_size=14, color=GREEN,
             alignment=PP_ALIGN.CENTER)

# Breakdown
add_shape_bg(slide, Inches(5.8), Inches(1.5), Inches(6.5), Inches(3.5), WHITE)
add_text_box(slide, Inches(6.1), Inches(1.7), Inches(6), Inches(0.4),
             'Breakdown', font_size=20, bold=True, color=DARK_BLUE)

breakdown_items = [
    'Criterion 1 (tickets \u2265 3):           73 participants',
    'Criterion 2 (surveys \u2265 2):          284 participants',
    'Criterion 3 (completed > rest):    64 participants',
    'Naive sum (if double-counted):   421',
    'Overlap removed:                           118',
    'Set union (actual):                         303  \u2705',
]
add_bullet_frame(slide, Inches(6.1), Inches(2.2), Inches(6), Inches(2.5),
                 breakdown_items, font_size=14)

# Proof
add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5), WHITE)
add_text_box(slide, Inches(1.1), Inches(5.5), Inches(11), Inches(0.4),
             'Inclusion-Exclusion Proof', font_size=18, bold=True, color=DARK_BLUE)
add_text_box(slide, Inches(1.1), Inches(5.9), Inches(11), Inches(0.7),
             '|A \u222a B \u222a C| = |A| + |B| + |C| \u2212 |A\u2229B| \u2212 |A\u2229C| \u2212 |B\u2229C| + |A\u2229B\u2229C|\n'
             '                 = 73 + 284 + 64 \u2212 57 \u2212 28 \u2212 61 + 28 = 303  \u2714',
             font_size=16, color=DARK_GRAY, font_name="Consolas")

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             'Tech Stack & Data Flow', font_size=36, bold=True, color=DARK_BLUE)

components = [
    ("Collibra API", "REST API\nBasic Auth", "Business terms,\nmeasures, column\nmappings, relations", ACCENT_BLUE),
    ("Delta Sharing", "Databricks\nOpen Protocol", "8 tables, ~20K rows\nSecure, no data copy", GREEN),
    ("Python / Pandas", "Computation\nEngine", "Data loading, SQL\nexecution, set logic", ORANGE),
    ("Claude (LLM)", "Anthropic\nclaude-sonnet-4-20250514", "SQL generation from\nnatural language +\ngoverned context", RGBColor(0x8E, 0x44, 0xAD)),
]

x_pos = Inches(0.3)
for title, subtitle, desc, color in components:
    add_shape_bg(slide, x_pos, Inches(1.5), Inches(3), Inches(2.5), WHITE)
    add_shape_bg(slide, x_pos, Inches(1.5), Inches(3), Inches(0.08), color)
    add_text_box(slide, x_pos + Inches(0.2), Inches(1.8), Inches(2.6), Inches(0.4),
                 title, font_size=20, bold=True, color=color)
    add_text_box(slide, x_pos + Inches(0.2), Inches(2.2), Inches(2.6), Inches(0.5),
                 subtitle, font_size=12, color=MID_GRAY)
    add_text_box(slide, x_pos + Inches(0.2), Inches(2.9), Inches(2.6), Inches(0.9),
                 desc, font_size=13, color=DARK_GRAY)
    x_pos += Inches(3.25)

# File list
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
             'Key Files', font_size=20, bold=True, color=DARK_BLUE)

file_items = [
    'collibra_client.py \u2014 Extracts all governed context from Collibra REST API',
    'data_loader.py \u2014 Loads Databricks tables via Delta Sharing into pandas',
    'ask_your_data.py \u2014 RAG pipeline: Collibra context + LLM + SQL execution',
    'rulebook.py \u2014 Governed definitions encoded as Python constants',
    'deterministic_engine.py \u2014 Pure pandas Active Tester computation (no LLM)',
    'test_deterministic.py \u2014 20-run verification harness',
]
add_bullet_frame(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(2.5),
                 file_items, font_size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             'Key Takeaways', font_size=36, bold=True, color=WHITE)

takeaways = [
    ("Context is not optional",
     "Without Collibra\u2019s governed definitions, the LLM returned 0 (wrong). With context, it returned 303 (correct). Same model, same data \u2014 only the context changed."),
    ("Governance enables determinism",
     "By encoding Collibra definitions as executable code, we achieve 100% reproducibility: 20 runs, 20 identical answers, zero LLM variability."),
    ("Two engines, one truth",
     "RAG for flexibility (handle any question), Rulebook for reliability (governed metrics). Both draw from the same Collibra source of truth."),
    ("Full traceability",
     "Every answer traces back to specific Collibra asset IDs, column mappings, and business definitions. Auditable, defensible, transparent."),
]

y_pos = Inches(1.6)
for title, desc in takeaways:
    add_shape_bg(slide, Inches(0.8), y_pos, Inches(11.5), Inches(1.2), RGBColor(0x22, 0x4A, 0x6E))
    add_text_box(slide, Inches(1.1), y_pos + Inches(0.1), Inches(11), Inches(0.4),
                 title, font_size=20, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.1), y_pos + Inches(0.5), Inches(11), Inches(0.6),
                 desc, font_size=14, color=RGBColor(0xD5, 0xDB, 0xDB))
    y_pos += Inches(1.4)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Thank You
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
             'Thank You', font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             'Ask Your Data \u2014 Governed AI for Enterprise Data Questions',
             font_size=22, color=RGBColor(0xAE, 0xD6, 0xF1), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
             'Team 18  |  Next Level Challenge Hackathon',
             font_size=18, color=RGBColor(0x85, 0x92, 0x9E), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "Ask_Your_Data_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
